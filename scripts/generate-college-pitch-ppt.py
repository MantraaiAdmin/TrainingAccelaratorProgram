#!/usr/bin/env python3
"""Generate Constel AI NextGen college pitch deck (.pptx)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "Constel_AI_NextGen_College_Pitch.pptx"

# Brand palette
INDIGO = RGBColor(0x31, 0x41, 0xB7)
INDIGO_DARK = RGBColor(0x1E, 0x1B, 0x4B)
SLATE = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x06, 0xB6, 0xD4)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left=Inches(0), top=Inches(0), width=Inches(0.12), height=Inches(7.5)):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def style_title(text_frame, size=32, color=INDIGO_DARK, bold=True):
    text_frame.word_wrap = True
    for i, p in enumerate(text_frame.paragraphs):
        p.font.size = Pt(size if i == 0 else size - 6)
        p.font.bold = bold if i == 0 else False
        p.font.color.rgb = color if i == 0 else SLATE
        p.font.name = "Calibri"
        p.space_after = Pt(8)


def add_title_slide(prs, title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, INDIGO_DARK)
    add_accent_bar(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5))

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(8.5), Inches(1.2))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"

    sub = slide.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(8.5), Inches(0.8))
    stf = sub.text_frame
    stf.text = subtitle
    sp = stf.paragraphs[0]
    sp.font.size = Pt(22)
    sp.font.color.rgb = ACCENT
    sp.font.name = "Calibri"

    foot = slide.shapes.add_textbox(Inches(0.9), Inches(6.2), Inches(8.5), Inches(0.6))
    ftf = foot.text_frame
    ftf.text = footer
    fp = ftf.paragraphs[0]
    fp.font.size = Pt(14)
    fp.font.color.rgb = MUTED
    fp.font.name = "Calibri"

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "2-min elevator: Constel AI NextGen is an internship platform built for engineering colleges. "
        "Eight weeks of industry-grade Python, full-stack, or AI — live coding, weekly assessments, "
        "interview prep, AI tutoring on syllabus. TPO gets one dashboard: bulk enroll, track completion, revenue share."
    )


def add_section_slide(prs, title, bullets, notes_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_accent_bar(slide)

    hdr = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(9), Inches(0.7))
    htf = hdr.text_frame
    htf.text = title
    style_title(htf, size=28, color=INDIGO_DARK)

    body = slide.shapes.add_textbox(Inches(0.65), Inches(1.35), Inches(8.8), Inches(5.5))
    btf = body.text_frame
    btf.word_wrap = True
    for i, item in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = SLATE
        p.font.name = "Calibri"
        p.space_after = Pt(10)
        if item.startswith("**") or item.startswith("•"):
            p.font.bold = item.startswith("**")

    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def add_table_slide(prs, title, headers, rows, notes_text="", col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_accent_bar(slide)

    hdr = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(9), Inches(0.7))
    htf = hdr.text_frame
    htf.text = title
    style_title(htf, size=26, color=INDIGO_DARK)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    left, top = Inches(0.55), Inches(1.25)
    width, height = Inches(9.0), Inches(0.45 * n_rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for ci, w in enumerate(col_widths):
            table.columns[ci].width = w

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = INDIGO
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFF)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = SLATE
                p.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, notes_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_accent_bar(slide)

    hdr = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(9), Inches(0.7))
    htf = hdr.text_frame
    htf.text = title
    style_title(htf, size=26, color=INDIGO_DARK)

    for col, (ctitle, items, x) in enumerate(
        [(left_title, left_items, 0.55), (right_title, right_items, 5.05)]
    ):
        tbox = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(4.2), Inches(0.5))
        ttf = tbox.text_frame
        ttf.text = ctitle
        tp = ttf.paragraphs[0]
        tp.font.bold = True
        tp.font.size = Pt(16)
        tp.font.color.rgb = INDIGO
        tp.font.name = "Calibri"

        bbox = slide.shapes.add_textbox(Inches(x), Inches(1.75), Inches(4.2), Inches(5.0))
        btf = bbox.text_frame
        btf.word_wrap = True
        for i, item in enumerate(items):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = SLATE
            p.font.name = "Calibri"
            p.space_after = Pt(8)

    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Constel AI NextGen",
        "Industry-driven AI education for tomorrow's innovators.",
        "Constel Global India Pvt. Ltd. · Praveen Manoharan",
    )

    add_section_slide(
        prs,
        "Company Background & Vision",
        [
            "13+ years of IT industry experience — we founded our startup to bridge academic learning and industry expectations",
            "Deep technical expertise serving clients across the United States",
            "We understand the skills and competencies organizations seek in today's workforce",
            "",
            "Mission: Transform students into job-ready professionals through industry-partnered internship programs",
            "",
            "Our Impact: A talent pipeline of industry-ready graduates who contribute confidently from day one",
        ],
        notes_text=(
            "Founders bring 13+ years IT experience and US client delivery. "
            "We built AI NextGen to close the academia-industry gap at scale with colleges."
        ),
    )

    add_section_slide(
        prs,
        "The Problem We Identified",
        [
            "Significant disconnect between skills from traditional education and what employers demand",
            "Graduates have theoretical knowledge but lack hands-on experience and industry exposure",
            "Limited professional readiness for competitive work environments",
            "Colleges face rising placement pressure with rapidly changing technology",
            "",
            "Bottom line: Students need measurable, industry-aligned internship pathways — not theory alone",
        ],
    )

    add_table_slide(
        prs,
        "The Employability Gap — By Stakeholder",
        ["Challenge", "Impact"],
        [
            ("Theory-heavy curriculum", "Students struggle in TCS/Infosys/product interviews"),
            ("No unified internship platform", "TPO tracks spreadsheets, WhatsApp, ad-hoc LMS"),
            ("Placement gap", "2nd–4th year lack Git, Linux, APIs, real projects"),
            ("Faculty bandwidth", "Cannot mentor 500+ students on industry tools"),
            ("Outdated lab exercises", "Output prediction ≠ system design ≠ AI engineering"),
        ],
        notes_text="Bottom line: Colleges need a scalable, measurable, industry-aligned internship layer — not another PDF syllabus.",
        col_widths=[Inches(3.2), Inches(5.8)],
    )

    add_section_slide(
        prs,
        "What We Do",
        [
            "Industry-partnered internship program designed to transform students into job-ready professionals",
            "Real-world project exposure · mentorship from experienced IT professionals",
            "Practical technical training · guidance on industry best practices",
            "",
            "Delivered via AI NextGen platform — full-stack learning & college management portal:",
            "Structured 8-week intensive tracks (~6 months of engineering depth)",
            "Live coding, quizzes, labs, capstones, interview prep — one portal",
            "Weekly live mentor sessions — student Q&A; recordings shared after each session",
            "Lifetime access · 24/7 lesson-scoped AI tutor · QR-verified credentials",
            "",
            "Not a video library. A job-readiness operating system for colleges.",
        ],
    )

    add_section_slide(
        prs,
        "Who It Is For",
        [
            "Students: 2nd–4th year CSE, IT, ECE, AI/ML, MCA",
            "College: TPO, department heads, internship coordinators",
            "Outcome: Git workflow, debug code, build APIs, face technical + HR interviews",
            "Ideal cohort: 50–500 students per department per semester",
        ],
    )

    add_section_slide(
        prs,
        "Student Journey",
        [
            "Enroll → Pick Track → Weekly Modules → Lessons + Labs → 50-Q Quiz (80% pass)",
            "→ Mini Projects → Interview Prep → Capstone → Certificate → Placement-ready portfolio",
            "",
            "Weekly rhythm:",
            "12–15 deep-dive topics with CLI, diagrams, industry scenarios",
            "Hands-on labs in Monaco editor (run code in browser)",
            "Weekly live session — address student queries; recorded sessions shared",
            "AI assistant for lesson-specific doubts (disabled during quizzes)",
            "Weekly assessment — 50 questions (MCQ, debugging, output prediction, scenarios)",
            "Interview module — industry Q&A, mock dialogue, unpredictable-question playbook",
        ],
    )

    add_section_slide(
        prs,
        "Three Master Tracks (Live Today)",
        [
            "Track 1 — Python Engineering Foundations (Beginner · 8 weeks)",
            "   Linux, Git, Python, DSA, OOP, REST APIs, capstone",
            "   Outcome: Backend-ready intern — automation, APIs, production debugging",
            "",
            "Track 2 — Full Stack Product Engineering (Intermediate · 8 weeks)",
            "   React, Next.js, NestJS, PostgreSQL, Docker, DevOps, real-time systems",
            "   Outcome: Junior full-stack engineer — ship features end-to-end",
            "",
            "Track 3 — AI Engineering & Intelligent Systems (Advanced · 8 weeks)",
            "   LLMs, RAG, agents, FastAPI, vector DBs, AI UX, production AI ops",
            "   Outcome: AI feature engineer — assistants, RAG pipelines, agents",
            "",
            "Each track: ~200+ lessons, coding exercises, interview prep, capstone options.",
        ],
    )

    add_table_slide(
        prs,
        "Curriculum Depth — Why It Feels Like 6 Months",
        ["Component", "What students get"],
        [
            ("Lessons", "8–10K chars each — CLI, walkthroughs, workflow diagrams, comparison tables"),
            ("Labs", "Monaco editor, test cases, hints, AI help"),
            ("Quizzes", "50 questions/week, 80% pass, anti-cheat (no AI during quiz)"),
            ("Live sessions", "Weekly live Q&A with mentors — recordings available on platform"),
            ("Interview prep", "8 industry Q&A + CLARIFY method + mock dialogue per module"),
            ("Capstone", "GitHub repo, deployment, architecture presentation"),
        ],
        notes_text="Example Week 1 (Python): Linux, terminal, Git, SSH, venv, debugging — 14 topics, 4 labs, mini-project.",
        col_widths=[Inches(2.2), Inches(6.8)],
    )

    add_section_slide(
        prs,
        "Platform Demo Highlights (5 min live)",
        [
            "Student dashboard — XP, progress, streak, assigned tracks",
            "Learn page — curriculum sidebar, lesson with diagrams, AI panel",
            "Coding lab — run Python, submit, get feedback",
            "Quiz — timed assessment, instant score",
            "Leaderboard & certificates",
            "Admin portal — add students, bulk upload, analytics, commission dashboard",
            "",
            "Demo login: student@demo.com / Demo@123",
            "Admin: admin@constel.ai / Demo@123",
        ],
        notes_text="Show live: dashboard, learn page with diagram + AI, coding lab, quiz, admin analytics.",
    )

    add_table_slide(
        prs,
        "AI That Colleges Can Trust",
        ["Feature", "Benefit"],
        [
            ("Lesson-scoped AI tutor", "Answers only current module — no life advice, no cheating"),
            ("Disabled during quizzes", "Academic integrity preserved"),
            ("Qwen / enterprise AI stack", "Cost-controlled, fallback providers"),
            ("Hints on coding exercises", "Guides without giving full solutions"),
        ],
        notes_text="Message to faculty: AI augments mentors — it does not replace evaluation.",
        col_widths=[Inches(3.5), Inches(5.5)],
    )

    add_table_slide(
        prs,
        "Why Constel vs Alternatives",
        ["Capability", "Constel", "Generic LMS", "YouTube / Udemy"],
        [
            ("Industry interview prep", "Built-in per week", "Rare", "None"),
            ("Live coding + sandbox", "Yes", "Rare", "No"),
            ("College admin + bulk enroll", "Yes", "Partial", "No"),
            ("Revenue share for college", "Yes", "No", "No"),
            ("AI tutor in context", "Yes", "Generic chatbot", "No"),
            ("Certificates + progress proof", "QR verified", "Varies", "No"),
            ("Indian placement season", "Built for it", "Generic", "Generic"),
        ],
        col_widths=[Inches(2.4), Inches(2.2), Inches(2.2), Inches(2.2)],
    )

    add_two_column_slide(
        prs,
        "College Partnership Model",
        "What Constel provides",
        [
            "Platform access (hosted or college-branded)",
            "Master curriculum + continuous updates",
            "Admin training for TPO/coordinators",
            "Student onboarding (bulk Excel upload)",
            "Analytics — completion, quiz scores, department progress",
            "Optional: mentor network, placement tie-ups",
        ],
        "What college provides",
        [
            "Cohort of students (2nd–4th year)",
            "TPO champion + 1 faculty liaison",
            "Lab / internet access (any laptop — cloud-based)",
            "Promotion within department",
        ],
        notes_text="Per-student investment by year. Includes lifetime access, live sessions, certificate, hackathon, LoR, and paid internship (₹5,000–10,000/mo stipend for top performers).",
    )

    add_section_slide(
        prs,
        "Outcomes & Placement Impact",
        [
            "After one 8-week track, students can demonstrate:",
            "Git workflow + PR culture",
            "Terminal/Linux comfort",
            "Working GitHub portfolio (labs + capstone)",
            "Weekly quiz scores as proof of mastery",
            "Interview prep modules completed",
            "Certificate with verification",
            "",
            "Career opportunities for top performers:",
            "Constel Hackathon — end-of-program competition across partner colleges",
            "Merit-based internship Letter of Recommendation (LoR)",
            "Paid internship at Constel — performance-based stipend ₹5,000–10,000/month",
            "",
            "TPO metrics: % completion by department, average quiz score,",
            "capstone submission rate, leaderboard top performers for recruiter shortlists",
        ],
    )

    add_table_slide(
        prs,
        "Implementation Timeline",
        ["Phase", "Duration", "Activity"],
        [
            ("1. MoU & setup", "Week 1", "Sign partnership, configure college, commission terms"),
            ("2. Pilot", "Week 2–3", "30–50 students, one track, faculty orientation"),
            ("3. Scale", "Week 4+", "Full department rollout, bulk upload"),
            ("4. Review", "End of Week 8", "Analytics review, certificate ceremony, placement prep"),
        ],
        notes_text="Technical requirement: Browser + internet. No on-premise servers required for pilot.",
        col_widths=[Inches(2.0), Inches(1.5), Inches(5.5)],
    )

    add_section_slide(
        prs,
        "Social Proof & Vision",
        [
            "Today:",
            "Production-ready platform (Next.js 15, NestJS, PostgreSQL, Docker sandbox)",
            "3 master tracks, 8 weeks each, 360+ structured submodules",
            "Admin analytics, finance dashboard, student CRUD",
            "",
            "Roadmap with college partners:",
            "Custom college branding · Razorpay billing · Faculty analytics",
            "Regional language support · Campus ambassador program",
            "",
            "Vision: Every Tier-2/3 college graduates students with",
            "GitHub + certificate + interview confidence — not just a degree.",
        ],
    )

    add_table_slide(
        prs,
        "Investment / Pricing",
        ["Student year", "Investment per student"],
        [
            ("2nd Year", "₹3,500"),
            ("3rd Year", "₹4,500"),
            ("4th Year", "₹6,999"),
        ],
        notes_text=(
            "Includes full 8-week track, lifetime application access, weekly live mentor sessions "
            "with recordings, certificate, hackathon, LoR eligibility, and paid internship "
            "(₹5,000–10,000/month performance-based stipend for top performers)."
        ),
        col_widths=[Inches(4.5), Inches(4.5)],
    )

    add_section_slide(
        prs,
        "Call to Action",
        [
            "Let's launch a 4-week pilot with [Department Name].",
            "",
            "Next steps:",
            "1. 30-minute live demo for TPO + HOD (today or this week)",
            "2. MoU draft + commission terms",
            "3. Bulk student upload (Excel template provided)",
            "4. Kickoff webinar for students",
            "5. Week-4 progress review → full semester contract",
            "",
            "Contact: Praveen Manoharan · praveen.manoharan@constelglobal.com · +91 8760 8760 62",
            "Constel Global India Pvt. Ltd. · https://www.constelglobal.com/",
        ],
    )

    # Appendix slides
    add_section_slide(
        prs,
        "Appendix A — Tech Stack",
        [
            "Next.js 15 · NestJS · PostgreSQL · Prisma",
            "Docker code sandbox · JWT auth · Swagger API",
            "For IT committee / technical review",
        ],
    )

    add_section_slide(
        prs,
        "Appendix B — Security",
        [
            "Role-based access (Student / Admin / Super Admin)",
            "Rate limiting · Secure code execution",
            "No secrets in client",
        ],
    )

    add_section_slide(
        prs,
        "Appendix C & D — Quiz & Capstones",
        [
            "Weekly quiz: 15 MCQ · 10 debugging · 10 output prediction",
            "5 scenario · 5 code completion · 5 problem solving",
            "",
            "Capstone examples:",
            "Python: productivity platform, automation suite",
            "Full-stack: LMS, HRMS SaaS",
            "AI: document intelligence, AI tutor, workflow agent",
        ],
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
