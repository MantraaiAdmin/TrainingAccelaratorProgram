#!/usr/bin/env python3
"""Generate AI NextGen Industry Readiness Program — 8-slide college pitch deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "AI_NextGen_Industry_Readiness_Pitch.pptx"
LOGO = ROOT / "docs" / "assets" / "constel-logo.png"

NAVY = RGBColor(0x0B, 0x1A, 0x3A)
NAVY_MID = RGBColor(0x15, 0x2A, 0x52)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
CYAN_DARK = RGBColor(0x08, 0x91, 0xB2)
SLATE = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
SKY = RGBColor(0x7D, 0xD3, 0xFC)

PROGRAM = "AI NextGen Industry Readiness Internship Program"


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.14), Inches(7.5))
    s.fill.solid()
    s.fill.fore_color.rgb = CYAN
    s.line.fill.background()


def title_box(slide, text, y=0.42, size=28, color=NAVY, light=False):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(y), Inches(9), Inches(0.75))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = WHITE if light else color


def bullet_box(slide, items, x, y, w, h, size=13, color=SLATE, prefix="• "):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{prefix}{item}" if prefix else item
        p.font.size = Pt(size)
        p.font.name = "Calibri"
        p.font.color.rgb = color
        p.space_after = Pt(6)


def card(slide, x, y, w, h, heading, lines, fill=WHITE, border=CYAN_DARK):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = heading
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.name = "Calibri"
    p.font.color.rgb = NAVY
    for line in lines:
        bp = tf.add_paragraph()
        bp.text = line
        bp.font.size = Pt(10)
        bp.font.name = "Calibri"
        bp.font.color.rgb = SLATE
        bp.space_before = Pt(4)


def arrow_flow(slide, labels, y=6.0):
    n = len(labels)
    step = 8.6 / n
    for i, label in enumerate(labels):
        x = 0.7 + i * step
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(step - 0.15), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = CYAN_DARK if i == n - 1 else NAVY_MID
        box.line.fill.background()
        tf = box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        if i < n - 1:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + step - 0.12), Inches(y + 0.18), Inches(0.18), Inches(0.2))
            arr.fill.solid()
            arr.fill.fore_color.rgb = CYAN
            arr.line.fill.background()


def slide1_about(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    bar(slide)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.75), Inches(0.55), Inches(1.1), Inches(1.1))

    brand = slide.shapes.add_textbox(Inches(2.0), Inches(0.65), Inches(7.5), Inches(0.5))
    brand.text_frame.text = PROGRAM
    brand.text_frame.paragraphs[0].font.size = Pt(16)
    brand.text_frame.paragraphs[0].font.bold = True
    brand.text_frame.paragraphs[0].font.color.rgb = CYAN
    brand.text_frame.paragraphs[0].font.name = "Calibri"

    title_box(slide, "Company Background & Vision", y=1.45, size=32, color=WHITE, light=True)

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(2.15), Inches(8.8), Inches(0.45))
    sub.text_frame.text = "Transforming Students into Industry-Ready Professionals"
    sub.text_frame.paragraphs[0].font.size = Pt(15)
    sub.text_frame.paragraphs[0].font.color.rgb = SKY
    sub.text_frame.paragraphs[0].font.name = "Calibri"

    bullet_box(
        slide,
        [
            "13+ years of IT industry experience — founded our startup to bridge academic learning and real-world expectations",
            "Deep technical expertise serving clients across the United States",
            "We understand the skills and competencies organizations seek in today's workforce",
            "Mission: Industry-partnered programs that transform students into job-ready professionals",
        ],
        0.75,
        2.75,
        8.8,
        2.0,
        size=13,
        color=RGBColor(0xE2, 0xE8, 0xF0),
    )

    impact = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.85), Inches(8.7), Inches(0.95))
    impact.fill.solid()
    impact.fill.fore_color.rgb = NAVY_MID
    impact.line.color.rgb = CYAN
    itf = impact.text_frame
    itf.text = "Our Impact"
    itf.paragraphs[0].font.bold = True
    itf.paragraphs[0].font.size = Pt(12)
    itf.paragraphs[0].font.color.rgb = CYAN
    ip = itf.add_paragraph()
    ip.text = (
        "Working with students, colleges, and industry partners to build technical, problem-solving, "
        "and workplace skills — a talent pipeline that contributes confidently from day one."
    )
    ip.font.size = Pt(11)
    ip.font.color.rgb = WHITE
    ip.font.name = "Calibri"

    highlights = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.95), Inches(8.7), Inches(0.85))
    highlights.fill.solid()
    highlights.fill.fore_color.rgb = NAVY_MID
    highlights.line.color.rgb = CYAN
    tf = highlights.text_frame
    tf.text = "Program Highlights"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = CYAN
    p2 = tf.add_paragraph()
    p2.text = (
        "24/7 AI Tutor  ·  Weekly Live Mentorship  ·  Lifetime Course Access  ·  "
        "Mini Project + Capstone + Hackathon  ·  Top 5% Paid Internship  ·  Top 5% Placement Assurity"
    )
    p2.font.size = Pt(10)
    p2.font.color.rgb = WHITE
    p2.font.name = "Calibri"

    arrow_flow(slide, ["Academia", "Industry Program", "Job-Ready Talent", "Day-One Impact"], y=6.05)


def slide2_challenges(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT)
    bar(slide)
    title_box(slide, "The Problem We Identified")

    narrative = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.15), Inches(8.9), Inches(1.05))
    narrative.fill.solid()
    narrative.fill.fore_color.rgb = RGBColor(0xE0, 0xF7, 0xFA)
    narrative.line.color.rgb = CYAN_DARK
    ntf = narrative.text_frame
    ntf.word_wrap = True
    ntf.text = (
        "Through industry engagements and hiring, we observed a significant disconnect: graduates "
        "have theoretical knowledge but lack hands-on experience, industry exposure, and professional "
        "readiness for competitive work environments."
    )
    ntf.paragraphs[0].font.size = Pt(11)
    ntf.paragraphs[0].font.color.rgb = SLATE
    ntf.paragraphs[0].font.name = "Calibri"

    card(
        slide,
        0.55,
        2.35,
        4.35,
        3.65,
        "Student Challenges",
        [
            "Strong theory, limited practical exposure",
            "No real-world project experience",
            "Limited industry workflow understanding",
            "Communication & interview readiness gaps",
            "Hard transition from classroom to workplace",
        ],
    )
    card(
        slide,
        5.1,
        2.35,
        4.35,
        3.65,
        "College Challenges",
        [
            "Rising placement expectations from management",
            "Rapidly changing technology landscape",
            "Industry–academia skill mismatch",
            "Need for measurable employability outcomes",
            "Faculty bandwidth for 500+ students",
        ],
    )

    gap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(6.15), Inches(5.6), Inches(0.7))
    gap.fill.solid()
    gap.fill.fore_color.rgb = NAVY
    gap.line.fill.background()
    tf = gap.text_frame
    tf.text = "College Education  ←—— GAP ——→  Industry Expectations"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def slide3_framework(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT)
    bar(slide)
    title_box(slide, "What We Do")

    intro = slide.shapes.add_textbox(Inches(0.55), Inches(1.12), Inches(8.9), Inches(0.55))
    intro.text_frame.word_wrap = True
    intro.text_frame.text = (
        "Industry-partnered internship program: real-world projects, mentorship from experienced IT "
        "professionals, practical technical training, and industry best practices."
    )
    intro.text_frame.paragraphs[0].font.size = Pt(11)
    intro.text_frame.paragraphs[0].font.color.rgb = SLATE
    intro.text_frame.paragraphs[0].font.name = "Calibri"

    pillars = [
        ("1", "Technical Skill\nDevelopment", "Python · Full Stack · AI tracks\n200+ lessons · live coding labs"),
        ("2", "Hands-on Project\nExperience", "Mini projects · capstone · GitHub portfolio"),
        ("3", "Industry\nMentorship", "Weekly live sessions · recorded Q&A · mentor feedback"),
        ("4", "Career Readiness\nTraining", "Interview prep · HR modules · mock dialogues"),
        ("5", "Assessment &\nTracking", "50-Q weekly quizzes · XP · leaderboard · analytics"),
        ("6", "Placement\nPreparation", "Aptitude · coding · resume · employer connect"),
    ]
    positions = [(0.55, 1.72), (3.45, 1.72), (6.35, 1.72), (0.55, 3.92), (3.45, 3.92), (6.35, 3.92)]
    for (num, title, desc), (x, y) in zip(pillars, positions):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.45), Inches(y), Inches(0.45), Inches(0.45))
        c.fill.solid()
        c.fill.fore_color.rgb = CYAN_DARK
        c.line.fill.background()
        c.text_frame.text = num
        c.text_frame.paragraphs[0].font.size = Pt(14)
        c.text_frame.paragraphs[0].font.bold = True
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        card(slide, x, y + 0.35, 2.75, 1.55, title.replace("\n", " "), [desc], fill=WHITE)

    outcomes = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.95), Inches(8.9), Inches(1.05))
    outcomes.fill.solid()
    outcomes.fill.fore_color.rgb = NAVY
    outcomes.line.fill.background()
    tf = outcomes.text_frame
    tf.text = "Our Impact"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = CYAN
    op = tf.add_paragraph()
    op.text = (
        "Industry-ready graduates with practical knowledge, portfolio proof, and confidence — "
        "prepared to thrive in modern organizations from day one."
    )
    op.font.size = Pt(11)
    op.font.color.rgb = WHITE
    op.alignment = PP_ALIGN.CENTER


def slide4_experiential(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT)
    bar(slide)
    title_box(slide, "Learn by Building, Not Just Learning")

    phases = [
        ("Phase 1", "Foundation Skills", "Programming · Databases · APIs · Git · Linux · Debugging"),
        ("Phase 2", "Project Development", "Team projects · Agile · Code reviews · Real-world use cases"),
        ("Phase 3", "Industry Simulation", "Professional workflows · Documentation · Collaboration"),
        ("Phase 4", "Performance Evaluation", "Weekly assessments · Feedback · Progress dashboards"),
    ]
    for i, (phase, name, detail) in enumerate(phases):
        x = 0.55 + i * 2.35
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.55), Inches(2.1), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = CYAN if i < 3 else CYAN_DARK
        line.line.fill.background()
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.95), Inches(1.35), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = CYAN_DARK
        dot.line.fill.background()
        card(slide, x, 1.85, 2.1, 2.0, f"{phase}: {name}", [detail])
        if i < 3:
            arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.05), Inches(2.55), Inches(0.25), Inches(0.25))
            arr.fill.solid()
            arr.fill.fore_color.rgb = CYAN
            arr.line.fill.background()

    tracks = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(4.2), Inches(8.9), Inches(2.5))
    tracks.fill.solid()
    tracks.fill.fore_color.rgb = RGBColor(0xE0, 0xF7, 0xFA)
    tracks.line.color.rgb = CYAN_DARK
    tf = tracks.text_frame
    tf.text = "Three Live Internship Tracks (8 Weeks Each)"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = NAVY
    for line in [
        "Python Engineering Foundations — Linux, Git, DSA, REST APIs, backend readiness",
        "Full Stack Product Engineering — React, Next.js, NestJS, PostgreSQL, Docker, DevOps",
        "AI Engineering & Intelligent Systems (Advanced) — LLMs, RAG, agents, FastAPI, vector DBs, AI UX",
    ]:
        p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(11)
        p.font.color.rgb = SLATE


def slide5_industry_exposure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT)
    bar(slide)
    title_box(slide, "Real-World Experience That Matters")

    bullet_box(
        slide,
        [
            "Live project environments with browser-based coding sandbox",
            "Industry-standard tools: Git, Docker, REST APIs, cloud-ready stacks",
            "Mentor-guided assignments with weekly live sessions + recordings",
            "Professional development: documentation, presentations, teamwork",
            "Client-centric problem solving through capstone & hackathon",
            "24/7 lesson-scoped AI tutor for continuous learning support",
        ],
        0.55,
        1.25,
        4.5,
        4.5,
        size=13,
    )

    card(
        slide,
        5.35,
        1.25,
        4.1,
        4.5,
        "Student Deliverables",
        [
            "✓  Project Portfolio (GitHub repos)",
            "✓  Technical Documentation",
            "✓  Presentation & Demo Skills",
            "✓  QR-Verified Completion Certificate",
            "✓  Internship LoR (merit-based)",
            "✓  Hackathon participation certificate",
        ],
        fill=WHITE,
    )

    collab = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.85), Inches(8.9), Inches(0.95))
    collab.fill.solid()
    collab.fill.fore_color.rgb = NAVY
    collab.line.fill.background()
    tf = collab.text_frame
    tf.text = "End of Program:  Mini Project  →  Capstone Project  →  Industry Hackathon  →  Top Performer Selection"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def slide6_placement_training(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT)
    bar(slide)
    title_box(slide, "Preparing Students for Career Success")

    sections = [
        ("Technical Preparation", ["Aptitude training", "Coding practice & output prediction", "Problem solving drills", "Mock technical interviews"]),
        ("Professional Skills", ["Communication workshops", "Resume building", "LinkedIn optimization", "Group discussions", "HR interview preparation"]),
        ("Career Readiness", ["Personal branding", "Workplace etiquette", "Corporate communication", "Weekly interview prep modules (8 Q&A per week)"]),
    ]
    for i, (heading, items) in enumerate(sections):
        card(slide, 0.55 + i * 3.15, 1.25, 2.95, 4.6, heading, [f"• {x}" for x in items])

    funnel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(6.0), Inches(3.6), Inches(0.75))
    funnel.fill.solid()
    funnel.fill.fore_color.rgb = CYAN_DARK
    funnel.line.fill.background()
    tf = funnel.text_frame
    tf.text = "Career Readiness Funnel:  Learn → Practice → Assess → Interview → Place"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def slide7_outcomes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT)
    bar(slide)
    title_box(slide, "Driving Employment Opportunities")

    card(
        slide, 0.55, 1.2, 2.85, 2.3, "Employer Engagement",
        ["Hiring partner network", "Recruitment drives", "Industry referrals", "Talent showcases at hackathon"],
    )
    card(
        slide, 3.55, 1.2, 2.85, 2.3, "Student Benefits",
        ["Increased placement opportunities", "Industry visibility", "Strong professional profile", "Job readiness confidence"],
    )
    card(
        slide, 6.55, 1.2, 2.85, 2.3, "Top Performers",
        ["Top 5% → Paid internship (performance-based stipend)", "Top 5% → Placement assurity pathway", "Hackathon winners prioritized"],
    )

    metrics = slide.shapes.add_table(2, 4, Inches(0.55), Inches(3.75), Inches(8.9), Inches(1.1)).table
    headers = ["Students Trained", "Projects Completed", "Hiring Partners", "Placement Success Rate"]
    values = ["[___]+", "[___]+", "[___]+", "[___]%"]
    for ci, (h, v) in enumerate(zip(headers, values)):
        metrics.cell(0, ci).text = h
        metrics.cell(1, ci).text = v
        for ri in range(2):
            cell = metrics.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if ri == 0 else LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11 if ri == 0 else 16)
                p.font.bold = ri == 1
                p.font.color.rgb = WHITE if ri == 0 else NAVY
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    fee = slide.shapes.add_table(4, 2, Inches(0.55), Inches(5.15), Inches(4.2), Inches(1.55)).table
    fee.cell(0, 0).text = "Student Year"
    fee.cell(0, 1).text = "Course Fee"
    fees = [("2nd Year", "₹3,999"), ("3rd Year", "₹7,999"), ("4th Year", "₹11,999")]
    for ri, (yr, price) in enumerate(fees, start=1):
        fee.cell(ri, 0).text = yr
        fee.cell(ri, 1).text = price
    for ri in range(4):
        for ci in range(2):
            cell = fee.cell(ri, ci)
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CYAN_DARK
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = ri == 0 or ci == 1
                p.font.color.rgb = WHITE if ri == 0 else SLATE

    eco = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(5.15), Inches(4.45), Inches(1.55))
    eco.fill.solid()
    eco.fill.fore_color.rgb = NAVY
    eco.line.fill.background()
    tf = eco.text_frame
    tf.text = "Placement Ecosystem"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = CYAN
    p = tf.add_paragraph()
    p.text = "Students  →  Skills + Portfolio  →  Hackathon  →  Employers  →  Internship / Placement"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE


def slide8_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, NAVY)
    bar(slide)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(4.45), Inches(0.65), Inches(1.1), Inches(1.1))

    title_box(slide, "Let's Build Industry-Ready Talent Together", y=1.85, size=30, color=WHITE, light=True)

    thank = slide.shapes.add_textbox(Inches(0.9), Inches(2.85), Inches(8.2), Inches(0.5))
    thank.text_frame.text = "Thank You  ·  Questions & Discussion"
    thank.text_frame.paragraphs[0].font.size = Pt(18)
    thank.text_frame.paragraphs[0].font.color.rgb = CYAN
    thank.text_frame.paragraphs[0].font.bold = True
    thank.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.55), Inches(8.2), Inches(1.1))
    quote.fill.solid()
    quote.fill.fore_color.rgb = NAVY_MID
    quote.line.color.rgb = CYAN
    tf = quote.text_frame
    tf.text = (
        '"Empowering students with the skills, experience, and confidence required '
        'to thrive in today\'s competitive job market."'
    )
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    contact = slide.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(8.2), Inches(1.5))
    ctf = contact.text_frame
    ctf.text = "Contact"
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.size = Pt(14)
    ctf.paragraphs[0].font.color.rgb = CYAN
    ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for line in [
        "Praveen Manoharan — Program Lead",
        "+91 8760 8760 62",
        PROGRAM,
    ]:
        p = ctf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide1_about(prs)
    slide2_challenges(prs)
    slide3_framework(prs)
    slide4_experiential(prs)
    slide5_industry_exposure(prs)
    slide6_placement_training(prs)
    slide7_outcomes(prs)
    slide8_closing(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
